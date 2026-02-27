"""
app/pptx_utils.py
PPTX 파일에서 슬라이드별 텍스트·표·다이어그램(차트/SmartArt)을 추출하여 마크다운으로 변환하는 모듈.
- 그룹 도형(GROUP) 내부를 재귀적으로 평탄화 후 top/left 순으로 정렬해 수집 (ssine/pptx2md 방식 참고).
- 표는 [[TABLE]]...[[/TABLE]], 차트/다이어그램/SmartArt는 [[DIAGRAM]]...[[/DIAGRAM]] 구분자로 감싼다.
"""

import logging
from operator import attrgetter

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape

from app.extract_constants import wrap_table, wrap_diagram

logger = logging.getLogger(__name__)


def _flatten_shapes(shapes) -> list:
    """
    그룹을 재귀적으로 펼쳐 비그룹 shape만 반환. ssine/pptx2md의 ungroup_shapes 패턴.
    """
    result = []
    for shape in shapes:
        try:
            is_group = (
                getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP
                or isinstance(shape, GroupShape)
            )
            if is_group and hasattr(shape, "shapes"):
                result.extend(_flatten_shapes(shape.shapes))
            else:
                result.append(shape)
        except Exception as e:
            logger.warning("shape 평탄화 스킵: %s", e)
    return result


def _shape_to_table_md(shape) -> str:
    """shape.table 을 마크다운 테이블 문자열로 변환. 병합 셀은 빈 문자열로 처리."""
    try:
        tbl = shape.table
    except Exception as e:
        logger.warning("shape.table 접근 실패: %s", e)
        return ""
    if not tbl.rows:
        return ""
    rows_md: list[str] = []
    for row in tbl.rows:
        cells = []
        for cell in row.cells:
            if getattr(cell, "is_spanned", False):
                cells.append("")
            else:
                try:
                    text = (cell.text or "").strip()
                    if not text and getattr(cell, "text_frame", None):
                        text = " ".join(p.text or "" for p in cell.text_frame.paragraphs).strip()
                    text = text.replace("\n", " ").replace("\v", " ")
                except Exception:
                    text = ""
                cells.append(text)
        rows_md.append("| " + " | ".join(cells) + " |")
    if not rows_md:
        return ""
    n_cols = len(tbl.rows[0].cells)
    rows_md.insert(1, "| " + " | ".join(["---"] * n_cols) + " |")
    return "\n".join(rows_md)


def _shape_to_diagram_caption(shape) -> str:
    """차트/다이어그램용 캡션. chart_title 이 있으면 사용, 없으면 기본 문구."""
    try:
        if getattr(shape, "has_chart", False):
            chart = shape.chart
            if chart.has_title and chart.chart_title.text_frame.text.strip():
                return chart.chart_title.text_frame.text.strip()
    except Exception:
        pass
    return "차트/다이어그램"


def _is_graphic_frame_diagram(shape) -> bool:
    """GraphicFrame이면서 차트/표가 아닌 경우(SmartArt 등) True."""
    if type(shape).__name__ != "GraphicFrame":
        return False
    if getattr(shape, "has_table", False) and shape.has_table:
        return False
    if getattr(shape, "has_chart", False) and shape.has_chart:
        return False
    return True


def _collect_from_shapes(shapes, title_holder: list[str]) -> list[str]:
    """
    shapes(및 그룹 내부)에서 텍스트·표·차트·다이어그램을 수집.
    그룹은 평탄화한 뒤 top/left 순으로 정렬해 시각적 읽기 순서로 처리 (ssine/pptx2md 참고).
    title_holder[0] 에 제목 플레이스홀더 텍스트가 설정될 수 있음.
    반환: body_parts (마크다운 조각 리스트)
    """
    flat = _flatten_shapes(shapes)
    try:
        sorted_shapes = sorted(flat, key=attrgetter("top", "left"))
    except Exception as e:
        logger.warning("shape 정렬 실패, 원본 순서 사용: %s", e)
        sorted_shapes = flat

    body_parts: list[str] = []

    for shape in sorted_shapes:
        # 표 (PlaceholderGraphicFrame 포함; has_table 또는 shape_type TABLE)
        has_table = getattr(shape, "has_table", False) and shape.has_table
        if not has_table and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TABLE:
            has_table = True
        if has_table:
            try:
                table_md = _shape_to_table_md(shape)
                if table_md:
                    body_parts.append(wrap_table(table_md))
            except Exception as e:
                logger.warning("표 추출 실패(shape 건너뜀): %s", e)
            continue

        # 차트
        if getattr(shape, "has_chart", False) and shape.has_chart:
            caption = _shape_to_diagram_caption(shape)
            body_parts.append(wrap_diagram(caption))
            continue

        # SmartArt 등 (GraphicFrame이지만 표/차트 아님)
        if _is_graphic_frame_diagram(shape):
            body_parts.append(wrap_diagram("SmartArt/다이어그램"))
            continue

        # 텍스트
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue

        if getattr(shape, "is_placeholder", False) and shape.is_placeholder:
            ph_type = getattr(shape.placeholder_format, "type", None)
            if ph_type in (1, 3, 13, 15):  # CENTER_TITLE, TITLE
                title_holder[0] = shape.text_frame.text.strip()
                continue

        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            level = getattr(para, "level", 0)
            indent = "  " * level
            body_parts.append(f"{indent}- {text}")

    return body_parts


def pptx_to_markdown(pptx_path: str, out_meta: dict | None = None) -> str:
    """
    PPTX 파일의 모든 슬라이드에서 텍스트·표·차트·SmartArt를 추출하여 마크다운으로 반환.
    - 그룹 도형 내부도 재귀 탐색하여 내용 수집.
    - 표: [[TABLE]]...[[/TABLE]], 차트/SmartArt: [[DIAGRAM]]...[[/DIAGRAM]]
    - out_meta 가 주어지면 slide_count 를 채움.
    """
    prs = Presentation(pptx_path)
    result_parts: list[str] = []
    slides = list(prs.slides)

    for slide_num, slide in enumerate(slides, start=1):
        title_holder: list[str] = [""]
        body_parts = _collect_from_shapes(slide.shapes, title_holder)
        title_text = title_holder[0]

        slide_md = f"## 🖼 Slide {slide_num}"
        if title_text:
            slide_md += f": {title_text}"
        slide_md += "\n\n"

        if body_parts:
            slide_md += "\n\n".join(body_parts)
        else:
            slide_md += "_(내용 없음)_"

        result_parts.append(slide_md)
        result_parts.append("\n\n---\n\n")

    if out_meta is not None:
        out_meta["slide_count"] = len(slides)

    return "\n".join(result_parts)
